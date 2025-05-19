import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation
prs = Presentation()

# Define colors
PRIMARY_COLOR = RGBColor(37, 99, 235)  # Blue
SECONDARY_COLOR = RGBColor(79, 70, 229)  # Indigo
ACCENT_COLOR = RGBColor(139, 92, 246)  # Purple
TEXT_COLOR = RGBColor(31, 41, 55)  # Dark gray
LIGHT_TEXT = RGBColor(107, 114, 128)  # Light gray
SUCCESS_COLOR = RGBColor(16, 185, 129)  # Green

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "AI Docx Builder Bootcamp"
subtitle.text = "From Concepts to Prototypes: Mastering AI-Powered Document Processing"

# Introduction slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "The Document Processing Revolution"
content.text = "• Organizations are drowning in documents\n"
content.text += "• Manual processing is slow, error-prone, and expensive\n"
content.text += "• AI offers transformative capabilities for document workflows\n"
content.text += "• Teams need practical skills to implement AI document solutions\n"
content.text += "• The AI Docx Builder Bootcamp bridges this critical skills gap"

# The Challenge slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "The Challenge"
content.text = "• 80% of business information resides in unstructured documents\n"
content.text += "• Organizations process thousands of documents daily\n"
content.text += "• Traditional document processing is:\n"
content.text += "   - Labor-intensive\n"
content.text += "   - Error-prone\n"
content.text += "   - Slow\n"
content.text += "   - Expensive\n"
content.text += "• Teams lack practical AI document processing skills"

# The Solution slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "The Solution: AI Docx Builder Bootcamp"
content.text = "A comprehensive training program that equips teams with practical AI document processing skills:\n\n"
content.text += "• Foundations On-Demand Modules\n"
content.text += "• Live Office Hours with AI Experts\n"
content.text += "• Hands-On Builder Assignments\n"
content.text += "• Block Hour Support\n"
content.text += "• Real-World Implementation Guidance"

# Program Structure slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Program Structure"
content.text = "Foundations On-Demand Modules:\n\n"
content.text += "1. Introduction to AI for Document Processing & Collaboration\n"
content.text += "2. Document Classification & Data Extraction with AI\n"
content.text += "3. Document Summarization & Task Automation\n"
content.text += "4. Security, Compliance, and Governance in AI Workflows\n"
content.text += "5. Integrating with Productivity Tools"

# Module 1 slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Module 1: Introduction to AI for Document Processing"
content.text = "• AI in Document Processing: An Overview\n"
content.text += "• Key Technologies and Capabilities\n"
content.text += "• Business Value and Use Cases\n"
content.text += "• Getting Started with AI Document Tools\n\n"
content.text += "Participants will understand fundamental AI concepts and their application to document workflows."

# Module 2 slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Module 2: Document Classification & Data Extraction"
content.text = "• Document Classification Fundamentals\n"
content.text += "• Extraction Techniques for Structured Documents\n"
content.text += "• Working with Semi-Structured Content\n"
content.text += "• Validation and Quality Control\n\n"
content.text += "Participants will learn how to automatically categorize documents and extract valuable data."

# Module 3 slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Module 3: Document Summarization & Task Automation"
content.text = "• Document Summarization Techniques\n"
content.text += "• Extractive vs. Abstractive Approaches\n"
content.text += "• Task Identification and Automation\n"
content.text += "• Building Automated Workflows\n\n"
content.text += "Participants will master techniques for condensing documents and automating routine tasks."

# Module 4 slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Module 4: Security, Compliance, and Governance"
content.text = "• Security Fundamentals for AI Document Systems\n"
content.text += "• Regulatory Compliance Frameworks\n"
content.text += "• Data Privacy and Protection\n"
content.text += "• Governance Best Practices\n\n"
content.text += "Participants will understand how to implement AI document solutions securely and responsibly."

# Module 5 slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Module 5: Integrating with Productivity Tools"
content.text = "• Integration Patterns and Approaches\n"
content.text += "• Microsoft 365 Integration\n"
content.text += "• Google Workspace Integration\n"
content.text += "• Custom API Development\n\n"
content.text += "Participants will learn how to connect AI document capabilities with everyday productivity tools."

# Live Components slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Interactive Learning Components"
content.text = "• Live Office Hours\n"
content.text += "   - Regular sessions with AI document processing experts\n"
content.text += "   - Q&A and problem-solving for specific challenges\n\n"
content.text += "• Hands-On Builder Assignments\n"
content.text += "   - Practical implementation exercises\n"
content.text += "   - Personalized feedback and guidance\n\n"
content.text += "• Block Hour Support\n"
content.text += "   - Dedicated assistance for specific implementation questions\n"
content.text += "   - Direct access to technical experts"

# Benefits slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Benefits for Your Organization"
content.text = "• Accelerate document processing by 70-90%\n"
content.text += "• Reduce manual processing costs by 60-80%\n"
content.text += "• Improve accuracy by 30-50%\n"
content.text += "• Enhance compliance and auditability\n"
content.text += "• Redeploy staff to higher-value activities\n"
content.text += "• Improve customer and employee experience\n"
content.text += "• Gain competitive advantage through AI adoption"

# Use Cases slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Real-World Applications"
content.text = "• Financial Services\n"
content.text += "   - Automated loan processing\n"
content.text += "   - Intelligent invoice handling\n\n"
content.text += "• Healthcare\n"
content.text += "   - Medical records processing\n"
content.text += "   - Insurance verification\n\n"
content.text += "• Legal\n"
content.text += "   - Contract analysis\n"
content.text += "   - Regulatory compliance monitoring\n\n"
content.text += "• General Business\n"
content.text += "   - HR document processing\n"
content.text += "   - Knowledge management"

# Pricing slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Investment Options"
content.text = "Individual Tiers:\n"
content.text += "• Basic Tier (Self-Paced): $1,495\n"
content.text += "• Professional Tier (Guided): $2,995\n"
content.text += "• Expert Tier (Comprehensive): $4,995\n\n"
content.text += "Corporate Tiers:\n"
content.text += "• Team Tier (5-10 participants): $12,500\n"
content.text += "• Department Tier (11-25 participants): $25,000\n"
content.text += "• Enterprise Tier (26-50 participants): $45,000\n\n"
content.text += "Flexible payment options and add-on services available"

# ROI slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Return on Investment"
content.text = "For a typical organization processing 10,000 documents monthly:\n\n"
content.text += "• Current cost: $5-15 per document = $50,000-$150,000 monthly\n"
content.text += "• With AI automation: $1-3 per document = $10,000-$30,000 monthly\n"
content.text += "• Monthly savings: $40,000-$120,000\n"
content.text += "• Annual savings: $480,000-$1,440,000\n\n"
content.text += "Investment in Department Tier ($25,000) pays for itself in less than one month"

# Competitive Advantage slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Why Choose AI Docx Builder Bootcamp?"
content.text = "• Document-Specific Focus\n"
content.text += "   - Specialized curriculum for document processing use cases\n\n"
content.text += "• Practical Implementation Skills\n"
content.text += "   - Hands-on experience with real-world scenarios\n\n"
content.text += "• Flexible Learning Options\n"
content.text += "   - Self-paced modules combined with live support\n\n"
content.text += "• Ongoing Expert Support\n"
content.text += "   - Guidance throughout your implementation journey\n\n"
content.text += "• Proven Methodology\n"
content.text += "   - Based on successful enterprise implementations"

# Implementation Timeline slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Implementation Timeline"
content.text = "Week 1-2: Foundations & Assessment\n"
content.text += "• Complete first two on-demand modules\n"
content.text += "• Assess current document workflows\n\n"
content.text += "Week 3-6: Skills Development\n"
content.text += "• Complete remaining modules\n"
content.text += "• Participate in live office hours\n"
content.text += "• Begin hands-on assignments\n\n"
content.text += "Week 7-8: Implementation Planning\n"
content.text += "• Develop organization-specific roadmap\n"
content.text += "• Create proof-of-concept implementations\n\n"
content.text += "Week 9-12: Deployment & Optimization\n"
content.text += "• Roll out initial AI document solutions\n"
content.text += "• Measure results and refine approach"

# Call to Action slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Next Steps"
content.text = "1. Schedule a detailed curriculum review\n\n"
content.text += "2. Identify key team members for initial training\n\n"
content.text += "3. Select appropriate tier based on organization size\n\n"
content.text += "4. Determine implementation timeline\n\n"
content.text += "5. Begin transforming your document workflows with AI"

# Contact slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Contact Information"
content.text = "For more information about the AI Docx Builder Bootcamp:\n\n"
content.text += "Email: ai.bootcamp@innovasolutions.com\n\n"
content.text += "Phone: (555) 123-4567\n\n"
content.text += "Website: www.innovasolutions.com/ai-bootcamp\n\n"
content.text += "Schedule a consultation: calendly.com/innova-ai-bootcamp"

# Save the presentation
output_path = os.path.join('/home/ubuntu/ai_docx_builder/presentations', 'AI_Docx_Builder_Bootcamp.pptx')
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)

print(f"Presentation saved to {output_path}")
