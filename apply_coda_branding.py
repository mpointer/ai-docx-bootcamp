import os
import re
from pathlib import Path

def update_html_files(directory):
    """Update HTML files with CODA branding and terminology"""
    html_files = list(Path(directory).glob('**/*.html'))
    
    for file_path in html_files:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Replace title and header references
        content = content.replace('AI Docx Builder Bootcamp', 'CODA Document Ecosystem')
        content = content.replace('From Concepts to Prototypes: Mastering AI-Powered Document Processing', 
                                 'Capture, Organize, Decide, Act: The Intelligent Document Ecosystem')
        
        # Update section titles and terminology
        content = content.replace('Foundations On-Demand Modules', 'CODA.core Components')
        content = content.replace('Office Hours', 'CODA.bins Implementation')
        content = content.replace('Hands-On Builder Assignments', 'CODA.vaults Configuration')
        content = content.replace('Block Hour Support', 'CODA.agents Deployment')
        
        # Update footer
        content = content.replace('&copy; 2025 AI Docx Builder Bootcamp', '&copy; 2025 CODA Document Ecosystem')
        
        with open(file_path, 'w', encoding='utf-8') as file:
            file.write(content)
        
        print(f"Updated HTML file: {file_path}")

def update_markdown_files(directory):
    """Update Markdown files with CODA branding and terminology"""
    md_files = list(Path(directory).glob('**/*.md'))
    
    for file_path in md_files:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Replace title and header references
        content = content.replace('AI Docx Builder Bootcamp', 'CODA Document Ecosystem')
        content = content.replace('# AI Docx Builder Bootcamp', '# CODA Document Ecosystem')
        
        # Update terminology
        content = content.replace('Bootcamp', 'Ecosystem')
        content = content.replace('bootcamp', 'ecosystem')
        content = content.replace('Foundations On-Demand modules', 'CODA.core components')
        content = content.replace('Office Hours', 'CODA.bins implementation')
        content = content.replace('Hands-On Builder Assignments', 'CODA.vaults configuration')
        content = content.replace('Block Hour Support', 'CODA.agents deployment')
        
        # Update specific pricing model references
        if 'pricing_model.md' in str(file_path):
            content = content.replace('Basic Tier: Self-Paced Learning', 'CODA.core: Ecosystem Foundation')
            content = content.replace('Professional Tier: Guided Learning', 'CODA.bins: Specialized Implementation')
            content = content.replace('Expert Tier: Comprehensive Learning', 'CODA.vaults + CODA.agents: Complete Solution')
            content = content.replace('Team Tier', 'Team Ecosystem')
            content = content.replace('Department Tier', 'Department Ecosystem')
            content = content.replace('Enterprise Tier', 'Enterprise Ecosystem')
        
        with open(file_path, 'w', encoding='utf-8') as file:
            file.write(content)
        
        print(f"Updated Markdown file: {file_path}")

def update_presentation(directory):
    """Create a script to update the PowerPoint presentation with CODA branding"""
    script_path = os.path.join(directory, 'update_coda_presentation.py')
    
    script_content = """
import os
from pptx import Presentation

# Path to the presentation
pptx_path = '/home/ubuntu/coda_update/presentations/AI_Docx_Builder_Bootcamp_updated.pptx'

# Load the presentation
prs = Presentation(pptx_path)

# Terms to replace
replacements = {
    'AI Docx Builder Bootcamp': 'CODA Document Ecosystem',
    'From Concepts to Prototypes': 'Capture, Organize, Decide, Act',
    'Bootcamp': 'Ecosystem',
    'bootcamp': 'ecosystem',
    'Foundations On-Demand': 'CODA.core',
    'Office Hours': 'CODA.bins',
    'Hands-On Builder Assignments': 'CODA.vaults',
    'Block Hour Support': 'CODA.agents',
    'Basic Tier': 'CODA.core',
    'Professional Tier': 'CODA.bins',
    'Expert Tier': 'CODA.vaults + CODA.agents',
    'Team Tier': 'Team Ecosystem',
    'Department Tier': 'Department Ecosystem',
    'Enterprise Tier': 'Enterprise Ecosystem'
}

# Process all slides
modified = False
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = original_text
                    
                    # Apply all replacements
                    for old_text, new_text_replacement in replacements.items():
                        if old_text in new_text:
                            new_text = new_text.replace(old_text, new_text_replacement)
                    
                    # Update if changed
                    if new_text != original_text:
                        run.text = new_text
                        modified = True

# Save the modified presentation
if modified:
    prs.save('/home/ubuntu/coda_update/presentations/CODA_Document_Ecosystem.pptx')
    print("Presentation updated and saved as CODA_Document_Ecosystem.pptx")
else:
    print("No changes made to the presentation")
"""
    
    with open(script_path, 'w', encoding='utf-8') as file:
        file.write(script_content)
    
    print(f"Created presentation update script: {script_path}")

def main():
    base_dir = '/home/ubuntu/coda_update'
    
    # Update HTML files
    update_html_files(base_dir)
    
    # Update Markdown files
    update_markdown_files(base_dir)
    
    # Create script to update PowerPoint
    update_presentation(base_dir)
    
    print("CODA branding updates completed. Run the presentation update script separately.")

if __name__ == "__main__":
    main()
