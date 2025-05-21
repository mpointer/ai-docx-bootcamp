
import os
from pptx import Presentation

# Path to the presentation
pptx_path = '/home/ubuntu/coda_update/presentations/AI_Docx_Builder_Bootcamp_updated.pptx'

# Load the presentation
prs = Presentation(pptx_path)

# Terms to replace
replacements = {
    'AI Docx Builder Bootcamp': 'CODA Document Ecosystem',
    'From Concepts to Prototypes': 'Capture, Organize, Decide, Act',
    'Bootcamp': 'Ecosystem',
    'bootcamp': 'ecosystem',
    'Foundations On-Demand': 'CODA.core',
    'Office Hours': 'CODA.bins',
    'Hands-On Builder Assignments': 'CODA.vaults',
    'Block Hour Support': 'CODA.agents',
    'Basic Tier': 'CODA.core',
    'Professional Tier': 'CODA.bins',
    'Expert Tier': 'CODA.vaults + CODA.agents',
    'Team Tier': 'Team Ecosystem',
    'Department Tier': 'Department Ecosystem',
    'Enterprise Tier': 'Enterprise Ecosystem'
}

# Process all slides
modified = False
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = original_text
                    
                    # Apply all replacements
                    for old_text, new_text_replacement in replacements.items():
                        if old_text in new_text:
                            new_text = new_text.replace(old_text, new_text_replacement)
                    
                    # Update if changed
                    if new_text != original_text:
                        run.text = new_text
                        modified = True

# Save the modified presentation
if modified:
    prs.save('/home/ubuntu/coda_update/presentations/CODA_Document_Ecosystem.pptx')
    print("Presentation updated and saved as CODA_Document_Ecosystem.pptx")
else:
    print("No changes made to the presentation")
