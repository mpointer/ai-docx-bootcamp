import os
from pptx import Presentation

# Path to the presentation
pptx_path = '/home/ubuntu/coda_update/presentations/AI_Docx_Builder_Bootcamp.pptx'

# Load the presentation
prs = Presentation(pptx_path)

# Words and phrases related to durations to look for and remove
duration_phrases = [
    "Duration:", 
    "minutes",
    "Week 1", 
    "Week 2", 
    "Week 3", 
    "Week 4", 
    "Week 5", 
    "Week 6",
    "Weeks 2-5",
    "weeks",
    "30 minutes",
    "45 minutes",
    "40 minutes",
    "35 minutes",
    "50 minutes"
]

# Function to check if text contains duration references
def contains_duration(text):
    if text is None:
        return False
    
    for phrase in duration_phrases:
        if phrase in text:
            return True
    return False

# Function to remove duration references from text
def remove_duration_references(text):
    if text is None:
        return None
    
    result = text
    for phrase in duration_phrases:
        if "Week" in phrase and phrase in result:
            # Replace week references with component names
            if "Week 1" in result:
                result = result.replace("Week 1", "Foundations")
            elif "Weeks 2-5" in result:
                result = result.replace("Weeks 2-5", "Support")
            elif "Week 6" in result:
                result = result.replace("Week 6", "Implementation")
        elif "Duration:" in result:
            # Remove the entire duration line
            parts = result.split('\n')
            filtered_parts = [part for part in parts if "Duration:" not in part]
            result = '\n'.join(filtered_parts)
        elif phrase in result:
            # Remove other duration references
            result = result.replace(phrase, "")
    
    return result

# Process all slides
modified = False
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if contains_duration(shape.text):
                new_text = remove_duration_references(shape.text)
                if shape.text != new_text:
                    # Update text frame with new text
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        # Clear existing paragraphs
                        p = text_frame.paragraphs[0]
                        p.clear()
                        p.text = new_text
                        modified = True

# Save the modified presentation
if modified:
    prs.save('/home/ubuntu/coda_update/presentations/AI_Docx_Builder_Bootcamp_updated.pptx')
    print("Presentation updated and saved as AI_Docx_Builder_Bootcamp_updated.pptx")
else:
    print("No duration references found in the presentation")
